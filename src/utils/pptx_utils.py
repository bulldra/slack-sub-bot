import logging
import tempfile
from typing import Iterable


def generate_pptx(slides: Iterable[str]) -> str:
    """Generate a PPTX file from slide texts.

    Returns the path to the created file or an empty string when the optional
    ``python-pptx`` library is unavailable.
    """
    try:
        from pptx import Presentation
    except Exception:  # pragma: no cover - optional dependency
        logging.warning("python-pptx is not installed")
        return ""

    prs = Presentation()
    layout = prs.slide_layouts[1]
    for text in slides:
        slide = prs.slides.add_slide(layout)
        placeholders = slide.shapes.placeholders
        if placeholders and len(placeholders) > 1:
            body = placeholders[1]
            body.text = str(text)
        else:
            # fallback when the layout does not have expected placeholders
            tx_box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
            tx_box.text = str(text)

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        return tmp.name
